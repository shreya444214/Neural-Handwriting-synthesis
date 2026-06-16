import os
import random
import string
import uuid
import re
from datetime import datetime, timedelta

from flask import Flask, request, jsonify, send_from_directory, send_file
from flask_cors import CORS
from flask_jwt_extended import (
    JWTManager, create_access_token, create_refresh_token,
    jwt_required, get_jwt_identity
)
from flask_socketio import SocketIO, emit, join_room, leave_room
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename

from config import Config
from models import db, User, File, ChatRoom, ChatMember, Message, LoginLog, OTP, HandwritingStyle

# ── App Setup ──────────────────────────────────────────────────────────────────

app = Flask(__name__)
app.config.from_object(Config)

CORS(app, resources={r"/api/*": {"origins": Config.CORS_ORIGINS}}, supports_credentials=True)
db.init_app(app)
jwt = JWTManager(app)
socketio = SocketIO(app, cors_allowed_origins=Config.CORS_ORIGINS, async_mode='threading')

# Create upload directories
os.makedirs(Config.UPLOAD_FOLDER, exist_ok=True)
os.makedirs(Config.PROFILE_PHOTOS_FOLDER, exist_ok=True)
os.makedirs(os.path.join(Config.UPLOAD_FOLDER, 'generated'), exist_ok=True)

with app.app_context():
    db.create_all()


# ── Health Check ───────────────────────────────────────────────────────────────

@app.route('/api/health', methods=['GET'])
def health_check():
    return jsonify({'status': 'ok', 'message': 'Handwriting Transformer API is running'}), 200

# ── Helpers ────────────────────────────────────────────────────────────────────

def generate_otp():
    return ''.join(random.choices(string.digits, k=6))

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in Config.ALLOWED_EXTENSIONS

def _setup_tesseract():
    """Configure Tesseract OCR path. Returns True if Tesseract is available."""
    try:
        import pytesseract
    except ImportError:
        return False, '[OCR requires pytesseract package. Install with: pip install pytesseract]'

    # Check common Windows install locations
    tesseract_paths = [
        r'C:\Program Files\Tesseract-OCR\tesseract.exe',
        r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
        r'C:\Users\HP\AppData\Local\Programs\Tesseract-OCR\tesseract.exe',
    ]
    for tp in tesseract_paths:
        if os.path.exists(tp):
            pytesseract.pytesseract.tesseract_cmd = tp
            return True, ''

    # Try system PATH
    import shutil
    if shutil.which('tesseract'):
        return True, ''

    return False, '[Tesseract-OCR not found. Please install it from https://github.com/UB-Mannheim/tesseract/wiki and ensure it is in your PATH or installed at C:\\Program Files\\Tesseract-OCR\\]'


def _find_poppler_path():
    """Find Poppler binaries for pdf2image. Returns path string or None."""
    import shutil
    # If pdftoppm is on PATH, no explicit path needed
    if shutil.which('pdftoppm'):
        return None

    # Check common install locations on Windows
    poppler_search = [
        r'C:\Program Files\poppler',
        r'C:\Program Files (x86)\poppler',
        r'C:\poppler',
        os.path.join(os.path.expanduser('~'), 'poppler'),
        os.path.join(os.path.expanduser('~'), 'Downloads', 'poppler'),
    ]
    for base in poppler_search:
        if os.path.isdir(base):
            # Search for bin directory containing pdftoppm
            for root, dirs, files in os.walk(base):
                if 'pdftoppm.exe' in files or 'pdftoppm' in files:
                    return root
            # Check common subdirectory patterns like poppler-xx.xx.x/Library/bin
            for item in os.listdir(base):
                bin_path = os.path.join(base, item, 'Library', 'bin')
                if os.path.isdir(bin_path):
                    return bin_path
                bin_path = os.path.join(base, item, 'bin')
                if os.path.isdir(bin_path):
                    return bin_path
    return None


def _preprocess_image_for_ocr(img, strategy='light'):
    """
    Apply preprocessing to a PIL Image for handwriting OCR.
    
    Strategies:
    - 'minimal': Just grayscale + resize. Best for clean, dark handwriting.
    - 'light': Grayscale + slight contrast boost + sharpen. Good general purpose.
    """
    from PIL import ImageFilter, ImageEnhance, ImageOps

    # Convert to RGB first if needed
    if img.mode != 'RGB':
        img = img.convert('RGB')

    # Convert to grayscale
    img = ImageOps.grayscale(img)

    # Resize to optimal OCR range (~3000px wide)
    # Tesseract works best at 300 DPI equivalent. Too large = slow + noisy.
    w, h = img.size
    target_w = 3000
    if w < 1500 or w > 4500:
        scale = target_w / w
        img = img.resize((int(w * scale), int(h * scale)), Image.LANCZOS)

    if strategy == 'minimal':
        return img

    if strategy == 'light':
        enhancer = ImageEnhance.Contrast(img)
        img = enhancer.enhance(1.5)
        img = img.filter(ImageFilter.SHARPEN)
        return img

    return img


def _score_ocr_quality(text):
    """
    Score OCR output quality. Higher = better readable English text.
    
    Measures:
    - Average word length (real words are 3-8 chars; garbled text has many 1-2 char fragments)
    - Ratio of real words (3+ letters) to total tokens
    - Clean character ratio (letters + spaces vs noise characters)
    
    Returns a float score where higher is better.
    """
    if not text or len(text.strip()) < 5:
        return 0.0

    words = text.split()
    if not words:
        return 0.0

    # Count words that look like real English (3+ letters, mostly alpha)
    real_words = 0
    total_word_len = 0
    for w in words:
        clean = ''.join(c for c in w if c.isalpha())
        if len(clean) >= 3:
            real_words += 1
            total_word_len += len(clean)

    if real_words == 0:
        return 0.0

    # Average length of real words (English averages ~4.7)
    avg_word_len = total_word_len / real_words

    # Ratio of real words to total tokens
    real_word_ratio = real_words / len(words)

    # Clean character ratio (letters + spaces + basic punctuation vs total)
    clean_chars = sum(1 for c in text if c.isalpha() or c in ' .,;:!?\n\'-')
    clean_ratio = clean_chars / len(text) if text else 0

    # Combined quality score
    # Reward: many real words, good average length, clean text
    score = (real_words * 2.0) * real_word_ratio * clean_ratio * min(avg_word_len / 4.0, 1.5)

    return score


def _ocr_extract_from_image(img):
    """
    Run Tesseract OCR on a PIL Image using multiple preprocessing strategies
    and PSM modes. Returns the best quality extracted text.
    
    Optimized: only 2 strategies × 2 PSM modes = 4 passes max.
    Uses quality-based scoring (word readability) rather than raw character count.
    """
    import pytesseract

    best_text = ''
    best_score = 0.0
    best_source = ''

    # Only 2 strategies — light works best for most handwriting, minimal as fallback
    strategies = ['light', 'minimal']

    # Only 2 PSM modes — 6 (uniform block) and 3 (fully automatic)
    psm_modes = [6, 3]

    for strategy in strategies:
        try:
            preprocessed = _preprocess_image_for_ocr(img, strategy=strategy)
        except Exception as e:
            print(f"    Preprocessing '{strategy}' failed: {e}")
            continue

        for psm in psm_modes:
            try:
                config = f'--oem 3 --psm {psm} -l eng'
                text = pytesseract.image_to_string(preprocessed, config=config).strip()

                score = _score_ocr_quality(text)

                if score > best_score:
                    best_text = text
                    best_score = score
                    best_source = f'{strategy}/psm{psm}'
                    alpha_count = sum(1 for c in text if c.isalpha())
                    print(f"    New best: {best_source} score={score:.1f} ({alpha_count} letters)")

                # Good enough — stop early
                if best_score > 30:
                    print(f"    Good quality achieved, stopping early")
                    return best_text

            except Exception:
                continue

    if best_source:
        alpha_count = sum(1 for c in best_text if c.isalpha())
        print(f"    Final best: {best_source} score={best_score:.1f}, {alpha_count} letters")

    return best_text


def extract_text_from_file(filepath, file_type):
    """Extract text content from uploaded files.
    
    For PDFs: first tries digital text extraction via PyPDF2.
    If that yields little/no text (handwritten/scanned PDF), falls back to
    OCR using pdf2image + pytesseract.
    """
    try:
        if file_type in ('txt',):
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()

        elif file_type in ('pdf',):
            # ── Step 1: Try digital text extraction (fast, works for typed PDFs) ──
            digital_text = ''
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(filepath)
                for page in reader.pages:
                    digital_text += page.extract_text() or ''
            except Exception as e:
                print(f"PyPDF2 extraction failed: {e}")

            # If we got reasonable text digitally, return it
            cleaned_digital = digital_text.strip()
            if len(cleaned_digital) > 30:
                print(f"PDF digital extraction succeeded: {len(cleaned_digital)} chars")
                return cleaned_digital

            # ── Step 2: Fall back to OCR for handwritten/scanned PDFs ──
            print(f"PDF digital text too short ({len(cleaned_digital)} chars), attempting OCR...")

            # Check Tesseract availability
            tess_ok, tess_err = _setup_tesseract()
            if not tess_ok:
                if cleaned_digital:
                    return cleaned_digital  # Return whatever PyPDF2 got
                return tess_err

            # Convert PDF pages to images using pdf2image
            try:
                from pdf2image import convert_from_path
                from PIL import Image
            except ImportError:
                if cleaned_digital:
                    return cleaned_digital
                return '[OCR for PDFs requires pdf2image package. Install with: pip install pdf2image. Also install Poppler: https://github.com/oschwartz10612/poppler-windows/releases]'

            try:
                poppler_path = _find_poppler_path()
                convert_kwargs = {'dpi': 300, 'fmt': 'png'}
                if poppler_path:
                    convert_kwargs['poppler_path'] = poppler_path

                pages = convert_from_path(filepath, **convert_kwargs)
                print(f"Converted PDF to {len(pages)} page image(s) for OCR at 400 DPI")
            except Exception as e:
                print(f"pdf2image conversion failed: {e}")
                import traceback
                traceback.print_exc()
                if cleaned_digital:
                    return cleaned_digital
                return f'[PDF to image conversion failed. Ensure Poppler is installed and in your PATH. Error: {str(e)}]'

            # OCR each page — _ocr_extract_from_image handles preprocessing internally
            import pytesseract
            all_text = []
            for i, page_img in enumerate(pages):
                try:
                    print(f"  Page {i+1}: running multi-strategy OCR...")
                    page_text = _ocr_extract_from_image(page_img)
                    if page_text:
                        all_text.append(page_text)
                        print(f"  Page {i+1}: extracted {len(page_text)} chars ({sum(1 for c in page_text if c.isalpha())} letters)")
                    else:
                        print(f"  Page {i+1}: no text extracted")
                except Exception as e:
                    print(f"  Page {i+1} OCR failed: {e}")

            ocr_result = '\n\n'.join(all_text).strip()

            if ocr_result:
                # Clean up common OCR artifacts
                ocr_result = _clean_ocr_text(ocr_result)
                print(f"PDF OCR complete: {len(ocr_result)} chars total")
                return ocr_result
            elif cleaned_digital:
                return cleaned_digital
            else:
                return '[No text could be extracted from this PDF. The document may be empty or the handwriting may be too unclear. Try a scan with higher contrast and darker ink.]'

        elif file_type in ('doc', 'docx'):
            try:
                from docx import Document
                doc = Document(filepath)
                return '\n'.join([p.text for p in doc.paragraphs])
            except Exception:
                return ''
        elif file_type in ('xls', 'xlsx'):
            try:
                from openpyxl import load_workbook
                wb = load_workbook(filepath, read_only=True)
                text = ''
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    for row in ws.iter_rows(values_only=True):
                        text += ' '.join([str(c) for c in row if c]) + '\n'
                return text
            except Exception:
                return ''
        elif file_type in ('ppt', 'pptx'):
            try:
                from pptx import Presentation
                prs = Presentation(filepath)
                text = ''
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            text += shape.text + '\n'
                return text
            except Exception:
                return ''
        elif file_type in ('jpg', 'jpeg', 'png', 'bmp', 'tiff', 'gif', 'webp'):
            try:
                from PIL import Image

                # Check Tesseract
                tess_ok, tess_err = _setup_tesseract()
                if not tess_ok:
                    return tess_err

                import pytesseract

                img = Image.open(filepath)
                extracted = _ocr_extract_from_image(img)

                if extracted:
                    return _clean_ocr_text(extracted)
                else:
                    return '[No text could be extracted from this image. Try a clearer image with darker handwriting on lighter paper.]'
            except Exception as e:
                print(f"OCR failed: {e}")
                import traceback
                traceback.print_exc()
                return f'[OCR extraction failed: {str(e)}]'
    except Exception:
        return ''
    return ''


def _clean_ocr_text(text):
    """Clean up common OCR artifacts from extracted text."""
    if not text:
        return text

    # Remove isolated single characters that are likely noise (except I, a, A)
    result = re.sub(r'(?<!\w)([^IaA\w\s])(?!\w)', '', text)

    # Fix multiple spaces
    result = re.sub(r' {2,}', ' ', result)

    # Fix multiple blank lines
    result = re.sub(r'\n{3,}', '\n\n', result)

    # Remove lines that are just punctuation/noise
    lines = result.split('\n')
    cleaned_lines = []
    for line in lines:
        stripped = line.strip()
        # Keep empty lines (paragraph breaks) and lines with actual words
        if not stripped:
            cleaned_lines.append('')
        elif len(re.findall(r'[a-zA-Z]', stripped)) >= 2:
            # Line has at least 2 letter characters — keep it
            cleaned_lines.append(line)
        elif len(stripped) > 3:
            # Line is long-ish — probably meaningful
            cleaned_lines.append(line)

    result = '\n'.join(cleaned_lines).strip()
    return result


def humanize_text(text):
    """Convert AI-generated text into more human/handwritten-style text."""
    if not text:
        return text
    
    # Simple humanization: add natural variations
    replacements = {
        'therefore': 'so',
        'however': 'but',
        'additionally': 'also',
        'furthermore': 'plus',
        'consequently': 'so',
        'nevertheless': 'still',
        'subsequently': 'then',
        'approximately': 'about',
        'utilize': 'use',
        'demonstrate': 'show',
        'implement': 'do',
        'facilitate': 'help',
        'endeavor': 'try',
        'commence': 'start',
        'terminate': 'end',
        'sufficient': 'enough',
        'numerous': 'many',
        'regarding': 'about',
        'in order to': 'to',
        'due to the fact that': 'because',
        'in the event that': 'if',
        'at this point in time': 'now',
        'for the purpose of': 'to',
        'in spite of the fact that': 'although',
        'it is important to note that': '',
        'it should be noted that': '',
    }
    
    result = text
    for formal, casual in replacements.items():
        result = re.sub(r'\b' + re.escape(formal) + r'\b', casual, result, flags=re.IGNORECASE)
    
    return result


def formalize_text(text):
    """Convert humanized/handwritten text into more formal AI-generated style.
    
    This function:
    1. Cleans up OCR artifacts (stray chars, broken words, bad spacing)
    2. Applies vocabulary formalization (casual → formal word replacements)
    3. Applies phrase-level upgrades (informal phrases → professional phrasing)
    4. Fixes sentence structure (capitalization, spacing, punctuation)
    """
    if not text:
        return text

    # ── Step 1: Clean up OCR artifacts ──
    result = _clean_ocr_text(text)

    # Fix common OCR misreads
    ocr_fixes = {
        r'\brn\b': 'm',         # "rn" often misread for "m"
        r'\bl\b(?=[a-z])': 'I', # lowercase L at start often should be I
        r'(?<=[a-z])0(?=[a-z])': 'o',  # zero in middle of word → letter o
        r'(?<=[a-z])1(?=[a-z])': 'l',  # one in middle of word → letter l
    }
    for pattern, fix in ocr_fixes.items():
        result = re.sub(pattern, fix, result)

    # ── Step 2: Word-level formalization ──
    word_replacements = {
        r'\bso\b': 'therefore',
        r'\bbut\b': 'however',
        r'\balso\b': 'additionally',
        r'\bplus\b': 'furthermore',
        r'\bstill\b': 'nevertheless',
        r'\bthen\b': 'subsequently',
        r'\babout\b': 'approximately',
        r'\buse\b': 'utilize',
        r'\bused\b': 'utilized',
        r'\buses\b': 'utilizes',
        r'\busing\b': 'utilizing',
        r'\bshow\b': 'demonstrate',
        r'\bshows\b': 'demonstrates',
        r'\bshowed\b': 'demonstrated',
        r'\bshowing\b': 'demonstrating',
        r'\bhelp\b': 'facilitate',
        r'\bhelps\b': 'facilitates',
        r'\bhelped\b': 'facilitated',
        r'\btry\b': 'endeavor',
        r'\btried\b': 'endeavored',
        r'\btrying\b': 'endeavoring',
        r'\bstart\b': 'commence',
        r'\bstarted\b': 'commenced',
        r'\bstarting\b': 'commencing',
        r'\bend\b': 'terminate',
        r'\bended\b': 'terminated',
        r'\bending\b': 'terminating',
        r'\benough\b': 'sufficient',
        r'\bmany\b': 'numerous',
        r'\bget\b': 'obtain',
        r'\bgets\b': 'obtains',
        r'\bgot\b': 'obtained',
        r'\bgetting\b': 'obtaining',
        r'\bgive\b': 'provide',
        r'\bgives\b': 'provides',
        r'\bgave\b': 'provided',
        r'\bgiving\b': 'providing',
        r'\bbig\b': 'significant',
        r'\bbigger\b': 'more significant',
        r'\bsmall\b': 'minimal',
        r'\bsmaller\b': 'more minimal',
        r'\bgood\b': 'excellent',
        r'\bbetter\b': 'superior',
        r'\bbest\b': 'optimal',
        r'\bbad\b': 'unfavorable',
        r'\bworse\b': 'more unfavorable',
        r'\bworst\b': 'most unfavorable',
        r'\blot of\b': 'substantial amount of',
        r'\blots of\b': 'substantial quantities of',
        r'\bkind of\b': 'somewhat',
        r'\bsort of\b': 'to a certain extent',
        r'\breally\b': 'substantially',
        r'\bvery\b': 'exceedingly',
        r'\bpretty\b': 'considerably',
        r'\bthing\b': 'element',
        r'\bthings\b': 'elements',
        r'\bstuff\b': 'material',
        r'\bway\b': 'methodology',
        r'\bways\b': 'methodologies',
        r'\bneed\b': 'require',
        r'\bneeds\b': 'requires',
        r'\bneeded\b': 'required',
        r'\bwant\b': 'desire',
        r'\bwants\b': 'desires',
        r'\bwanted\b': 'desired',
        r'\bmake\b': 'construct',
        r'\bmakes\b': 'constructs',
        r'\bmade\b': 'constructed',
        r'\bmaking\b': 'constructing',
        r'\bkeep\b': 'maintain',
        r'\bkeeps\b': 'maintains',
        r'\bkept\b': 'maintained',
        r'\bkeeping\b': 'maintaining',
        r'\blook\b': 'examine',
        r'\blooks\b': 'examines',
        r'\blooked\b': 'examined',
        r'\blooking\b': 'examining',
        r'\bfind\b': 'identify',
        r'\bfinds\b': 'identifies',
        r'\bfound\b': 'identified',
        r'\bfinding\b': 'identifying',
        r'\btalk\b': 'discuss',
        r'\btalks\b': 'discusses',
        r'\btalked\b': 'discussed',
        r'\btalking\b': 'discussing',
        r'\bask\b': 'inquire',
        r'\basks\b': 'inquires',
        r'\basked\b': 'inquired',
        r'\basking\b': 'inquiring',
        r'\bfix\b': 'rectify',
        r'\bfixed\b': 'rectified',
        r'\bfixing\b': 'rectifying',
        r'\bcheck\b': 'verify',
        r'\bchecks\b': 'verifies',
        r'\bchecked\b': 'verified',
        r'\bchecking\b': 'verifying',
        r'\bbuild\b': 'construct',
        r'\bbuilds\b': 'constructs',
        r'\bbuilt\b': 'constructed',
        r'\bbuilding\b': 'constructing',
        r'\bgo\b': 'proceed',
        r'\bgoes\b': 'proceeds',
        r'\bgoing\b': 'proceeding',
        r'\bgone\b': 'proceeded',
        r'\bput\b': 'place',
        r'\bputs\b': 'places',
        r'\bputting\b': 'placing',
        r'\btell\b': 'inform',
        r'\btells\b': 'informs',
        r'\btold\b': 'informed',
        r'\btelling\b': 'informing',
        r'\bthink\b': 'consider',
        r'\bthinks\b': 'considers',
        r'\bthought\b': 'considered',
        r'\bthinking\b': 'considering',
    }

    for pattern, replacement in word_replacements.items():
        result = re.sub(pattern, replacement, result, flags=re.IGNORECASE)

    # ── Step 3: Phrase-level upgrades ──
    phrase_replacements = [
        (r'\ba lot\b', 'a considerable amount'),
        (r'\bright now\b', 'at the present time'),
        (r'\bin the end\b', 'ultimately'),
        (r'\bbecause of\b', 'due to'),
        (r'\bcome up with\b', 'devise'),
        (r'\bpoint out\b', 'highlight'),
        (r'\bset up\b', 'establish'),
        (r'\bcarry out\b', 'execute'),
        (r'\bfigure out\b', 'determine'),
        (r'\blook into\b', 'investigate'),
        (r'\bbring up\b', 'introduce'),
        (r'\bgive up\b', 'relinquish'),
        (r'\bgo through\b', 'undergo'),
        (r'\bput together\b', 'assemble'),
        (r'\btake care of\b', 'manage'),
        (r'\bmake sure\b', 'ensure'),
        (r'\bfor example\b', 'for instance'),
        (r'\bin fact\b', 'indeed'),
        (r'\bas well\b', 'in addition'),
        (r'\bin short\b', 'in summary'),
    ]

    for pattern, replacement in phrase_replacements:
        result = re.sub(pattern, replacement, result, flags=re.IGNORECASE)

    # ── Step 4: Fix sentence structure ──
    # Fix double spaces
    result = re.sub(r' {2,}', ' ', result)

    # Ensure sentences end with proper punctuation
    lines = result.split('\n')
    formatted_lines = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            formatted_lines.append('')
            continue

        # Capitalize first letter of each sentence
        sentences = re.split(r'(?<=[.!?])\s+', stripped)
        capitalized = []
        for s in sentences:
            s = s.strip()
            if s:
                s = s[0].upper() + s[1:] if len(s) > 1 else s.upper()
                capitalized.append(s)
        line_result = ' '.join(capitalized)

        # Add period at end if the line doesn't end with punctuation
        if line_result and line_result[-1] not in '.!?:;,':
            line_result += '.'

        formatted_lines.append(line_result)

    result = '\n'.join(formatted_lines)

    # Fix spacing around punctuation
    result = re.sub(r'\s+([.,;:!?])', r'\1', result)
    result = re.sub(r'([.,;:!?])(?=[A-Za-z])', r'\1 ', result)

    # Remove duplicate periods
    result = re.sub(r'\.{2,}', '.', result)

    return result


# ── Auth Endpoints ─────────────────────────────────────────────────────────────

@app.route('/api/auth/send-otp', methods=['POST'])
def send_otp():
    data = request.get_json()
    target = data.get('target', '').strip()
    purpose = data.get('purpose', 'register')
    
    if not target:
        return jsonify({'error': 'Email or phone number is required'}), 400
    
    # Invalidate previous OTPs
    OTP.query.filter_by(target=target, is_used=False).update({'is_used': True})
    
    otp_code = generate_otp()
    new_otp = OTP(
        target=target,
        otp_code=otp_code,
        purpose=purpose,
        expires_at=datetime.utcnow() + timedelta(minutes=Config.OTP_EXPIRY_MINUTES)
    )
    db.session.add(new_otp)
    db.session.commit()
    
    # In production, send via email/SMS service
    print(f"\n{'='*50}")
    print(f"  OTP for {target}: {otp_code}")
    print(f"  Purpose: {purpose}")
    print(f"  Expires in {Config.OTP_EXPIRY_MINUTES} minutes")
    print(f"{'='*50}\n")
    
    return jsonify({
        'message': f'OTP sent to {target}',
        'otp_debug': otp_code  # Remove in production!
    }), 200


@app.route('/api/auth/verify-otp', methods=['POST'])
def verify_otp():
    data = request.get_json()
    target = data.get('target', '').strip()
    otp_code = data.get('otp', '').strip()
    purpose = data.get('purpose', 'register')
    
    if not target or not otp_code:
        return jsonify({'error': 'Target and OTP are required'}), 400
    
    otp_record = OTP.query.filter_by(
        target=target,
        otp_code=otp_code,
        purpose=purpose,
        is_used=False
    ).first()
    
    if not otp_record:
        return jsonify({'error': 'Invalid OTP'}), 400
    
    if otp_record.expires_at < datetime.utcnow():
        return jsonify({'error': 'OTP has expired'}), 400
    
    otp_record.is_used = True
    db.session.commit()
    
    return jsonify({'message': 'OTP verified successfully', 'verified': True}), 200


@app.route('/api/auth/register', methods=['POST'])
def register():
    data = request.get_json()
    
    name = data.get('name', '').strip()
    username = data.get('username', '').strip()
    email = data.get('email', '').strip().lower()
    phone = data.get('phone', '').strip()
    password = data.get('password', '')
    dob = data.get('dob', '')
    gender = data.get('gender', '')
    profession = data.get('profession', '')
    
    # Validation
    if not all([name, username, email, password]):
        return jsonify({'error': 'Name, username, email, and password are required'}), 400
    
    if len(password) < 6:
        return jsonify({'error': 'Password must be at least 6 characters'}), 400
    
    if User.query.filter_by(email=email).first():
        return jsonify({'error': 'Email already registered'}), 409
    
    if User.query.filter_by(username=username).first():
        return jsonify({'error': 'Username already taken'}), 409
    
    new_user = User(
        name=name,
        username=username,
        email=email,
        phone=phone,
        password_hash=generate_password_hash(password),
        dob=dob,
        gender=gender,
        profession=profession,
        is_verified=True,
        is_email_verified=True,
        is_phone_verified=bool(phone)
    )
    
    db.session.add(new_user)
    db.session.commit()
    
    access_token = create_access_token(identity=str(new_user.id))
    refresh_token = create_refresh_token(identity=str(new_user.id))
    
    # Log registration
    log = LoginLog(user_id=new_user.id, action='register', ip_address=request.remote_addr)
    db.session.add(log)
    db.session.commit()
    
    return jsonify({
        'message': f'Welcome {name}! Registration successful.',
        'access_token': access_token,
        'refresh_token': refresh_token,
        'user': new_user.to_dict()
    }), 201


@app.route('/api/auth/login', methods=['POST'])
def login():
    data = request.get_json()
    email = data.get('email', '').strip().lower()
    password = data.get('password', '')
    
    if not email or not password:
        return jsonify({'error': 'Email and password are required'}), 400
    
    user = User.query.filter_by(email=email).first()
    
    if not user or not check_password_hash(user.password_hash, password):
        return jsonify({'error': 'Invalid email or password'}), 401
    
    user.last_login = datetime.utcnow()
    
    access_token = create_access_token(identity=str(user.id))
    refresh_token = create_refresh_token(identity=str(user.id))
    
    # Log login
    log = LoginLog(user_id=user.id, action='login', ip_address=request.remote_addr,
                   user_agent=request.headers.get('User-Agent', '')[:256])
    db.session.add(log)
    db.session.commit()
    
    return jsonify({
        'message': f'Welcome back, {user.name}!',
        'access_token': access_token,
        'refresh_token': refresh_token,
        'user': user.to_dict()
    }), 200


@app.route('/api/auth/logout', methods=['POST'])
@jwt_required()
def logout():
    user_id = int(get_jwt_identity())
    log = LoginLog(user_id=user_id, action='logout', ip_address=request.remote_addr)
    db.session.add(log)
    db.session.commit()
    return jsonify({'message': 'Logged out successfully'}), 200


@app.route('/api/auth/forgot-password', methods=['POST'])
def forgot_password():
    data = request.get_json()
    email = data.get('email', '').strip().lower()
    
    if not email:
        return jsonify({'error': 'Email is required'}), 400
    
    user = User.query.filter_by(email=email).first()
    if not user:
        return jsonify({'error': 'No account found with this email'}), 404
    
    # Invalidate previous OTPs
    OTP.query.filter_by(target=email, purpose='reset_password', is_used=False).update({'is_used': True})
    
    otp_code = generate_otp()
    new_otp = OTP(
        target=email,
        otp_code=otp_code,
        purpose='reset_password',
        expires_at=datetime.utcnow() + timedelta(minutes=Config.OTP_EXPIRY_MINUTES)
    )
    db.session.add(new_otp)
    db.session.commit()
    
    print(f"\n{'='*50}")
    print(f"  PASSWORD RESET OTP for {email}: {otp_code}")
    print(f"{'='*50}\n")
    
    return jsonify({
        'message': 'Password reset OTP sent to your email',
        'otp_debug': otp_code
    }), 200


@app.route('/api/auth/reset-password', methods=['POST'])
def reset_password():
    data = request.get_json()
    email = data.get('email', '').strip().lower()
    otp_code = data.get('otp', '').strip()
    new_password = data.get('new_password', '')
    
    if not all([email, otp_code, new_password]):
        return jsonify({'error': 'Email, OTP, and new password are required'}), 400
    
    if len(new_password) < 6:
        return jsonify({'error': 'Password must be at least 6 characters'}), 400
    
    otp_record = OTP.query.filter_by(
        target=email, otp_code=otp_code, purpose='reset_password', is_used=False
    ).first()
    
    if not otp_record or otp_record.expires_at < datetime.utcnow():
        return jsonify({'error': 'Invalid or expired OTP'}), 400
    
    user = User.query.filter_by(email=email).first()
    if not user:
        return jsonify({'error': 'User not found'}), 404
    
    user.password_hash = generate_password_hash(new_password)
    otp_record.is_used = True
    db.session.commit()
    
    return jsonify({'message': 'Password reset successfully'}), 200


@app.route('/api/auth/me', methods=['GET'])
@jwt_required()
def get_me():
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    return jsonify({'user': user.to_dict()}), 200


@app.route('/api/auth/refresh', methods=['POST'])
@jwt_required(refresh=True)
def refresh_token():
    user_id = get_jwt_identity()
    new_token = create_access_token(identity=user_id)
    return jsonify({'access_token': new_token}), 200


# ── Profile Endpoints ──────────────────────────────────────────────────────────

@app.route('/api/profile', methods=['GET'])
@jwt_required()
def get_profile():
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    
    file_count = File.query.filter_by(user_id=user_id).count()
    
    return jsonify({
        'user': user.to_dict(),
        'stats': {
            'total_files': file_count,
            'account_age_days': (datetime.utcnow() - user.created_at).days if user.created_at else 0,
        }
    }), 200


@app.route('/api/profile', methods=['PUT'])
@jwt_required()
def update_profile():
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    
    data = request.get_json()
    
    if 'name' in data:
        user.name = data['name'].strip()
    if 'phone' in data:
        user.phone = data['phone'].strip()
    if 'dob' in data:
        user.dob = data['dob']
    if 'gender' in data:
        user.gender = data['gender']
    if 'profession' in data:
        user.profession = data['profession']
    if 'theme' in data:
        user.theme = data['theme']
    if 'notifications_enabled' in data:
        user.notifications_enabled = data['notifications_enabled']
    
    db.session.commit()
    return jsonify({'message': 'Profile updated successfully', 'user': user.to_dict()}), 200


@app.route('/api/profile/photo', methods=['POST'])
@jwt_required()
def upload_profile_photo():
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    
    if 'photo' not in request.files:
        return jsonify({'error': 'No photo file provided'}), 400
    
    file = request.files['photo']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    ext = file.filename.rsplit('.', 1)[1].lower() if '.' in file.filename else 'jpg'
    filename = f"profile_{user_id}_{uuid.uuid4().hex[:8]}.{ext}"
    filepath = os.path.join(Config.PROFILE_PHOTOS_FOLDER, filename)
    file.save(filepath)
    
    user.profile_photo = filename
    db.session.commit()
    
    return jsonify({'message': 'Profile photo updated', 'photo': filename}), 200


@app.route('/api/profile/photo/<filename>', methods=['GET'])
def serve_profile_photo(filename):
    return send_from_directory(Config.PROFILE_PHOTOS_FOLDER, filename)


@app.route('/api/profile/change-password', methods=['POST'])
@jwt_required()
def change_password():
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    
    data = request.get_json()
    current_password = data.get('current_password', '')
    new_password = data.get('new_password', '')
    
    if not check_password_hash(user.password_hash, current_password):
        return jsonify({'error': 'Current password is incorrect'}), 401
    
    if len(new_password) < 6:
        return jsonify({'error': 'New password must be at least 6 characters'}), 400
    
    user.password_hash = generate_password_hash(new_password)
    db.session.commit()
    
    return jsonify({'message': 'Password changed successfully'}), 200


# ── File Endpoints ─────────────────────────────────────────────────────────────

@app.route('/api/files/upload', methods=['POST'])
@jwt_required()
def upload_file():
    user_id = int(get_jwt_identity())
    
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not allowed_file(file.filename):
        return jsonify({'error': f'File type not allowed. Allowed: {", ".join(Config.ALLOWED_EXTENSIONS)}'}), 400
    
    original_name = secure_filename(file.filename)
    ext = original_name.rsplit('.', 1)[1].lower() if '.' in original_name else 'txt'
    stored_name = f"{uuid.uuid4().hex}.{ext}"
    filepath = os.path.join(Config.UPLOAD_FOLDER, stored_name)
    
    file.save(filepath)
    file_size = os.path.getsize(filepath)
    
    # Extract text
    content_text = extract_text_from_file(filepath, ext)
    
    new_file = File(
        user_id=user_id,
        original_name=file.filename,
        stored_name=stored_name,
        file_type=ext,
        file_size=file_size,
        status='uploaded',
        content_text=content_text
    )
    
    db.session.add(new_file)
    db.session.commit()
    
    return jsonify({
        'message': 'File uploaded successfully',
        'file': new_file.to_dict()
    }), 201


@app.route('/api/files/upload-text', methods=['POST'])
@jwt_required()
def upload_text():
    user_id = int(get_jwt_identity())
    data = request.get_json()
    
    text = data.get('text', '').strip()
    name = data.get('name', 'Pasted Text').strip()
    
    if not text:
        return jsonify({'error': 'Text content is required'}), 400
    
    stored_name = f"{uuid.uuid4().hex}.txt"
    filepath = os.path.join(Config.UPLOAD_FOLDER, stored_name)
    
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write(text)
    
    new_file = File(
        user_id=user_id,
        original_name=name,
        stored_name=stored_name,
        file_type='txt',
        file_size=len(text.encode('utf-8')),
        status='uploaded',
        content_text=text
    )
    
    db.session.add(new_file)
    db.session.commit()
    
    return jsonify({
        'message': 'Text saved successfully',
        'file': new_file.to_dict()
    }), 201


@app.route('/api/files', methods=['GET'])
@jwt_required()
def list_files():
    user_id = int(get_jwt_identity())
    files = File.query.filter_by(user_id=user_id).order_by(File.created_at.desc()).all()
    return jsonify({'files': [f.to_dict() for f in files]}), 200


@app.route('/api/files/<int:file_id>', methods=['GET'])
@jwt_required()
def get_file(file_id):
    user_id = int(get_jwt_identity())
    file = File.query.filter_by(id=file_id, user_id=user_id).first()
    if not file:
        return jsonify({'error': 'File not found'}), 404
    return jsonify({'file': file.to_dict()}), 200


@app.route('/api/files/<int:file_id>', methods=['DELETE'])
@jwt_required()
def delete_file(file_id):
    user_id = int(get_jwt_identity())
    file = File.query.filter_by(id=file_id, user_id=user_id).first()
    if not file:
        return jsonify({'error': 'File not found'}), 404
    
    # Delete physical file
    filepath = os.path.join(Config.UPLOAD_FOLDER, file.stored_name)
    if os.path.exists(filepath):
        os.remove(filepath)
    
    db.session.delete(file)
    db.session.commit()
    
    return jsonify({'message': 'File deleted successfully'}), 200


@app.route('/api/files/<int:file_id>/transform', methods=['POST'])
@jwt_required()
def transform_file(file_id):
    user_id = int(get_jwt_identity())
    file = File.query.filter_by(id=file_id, user_id=user_id).first()
    if not file:
        return jsonify({'error': 'File not found'}), 404
    
    data = request.get_json()
    transform_type = data.get('transform_type', 'humanized')  # humanized or ai_generated
    text_content = data.get('text_content', file.content_text or '')
    
    # Apply transformation
    if transform_type == 'humanized':
        transformed = humanize_text(text_content)
    elif transform_type == 'ai_generated':
        transformed = formalize_text(text_content)
    else:
        transformed = text_content
    
    file.status = 'completed'
    file.transform_type = transform_type
    file.transformed_text = transformed
    file.updated_at = datetime.utcnow()
    
    db.session.commit()
    
    return jsonify({
        'message': f'File transformed to {transform_type} successfully',
        'file': file.to_dict(),
        'transformed_text': transformed,
    }), 200


@app.route('/api/files/transform-text', methods=['POST'])
@jwt_required()
def transform_text_direct():
    """Transform text directly without a file — for quick transformations."""
    data = request.get_json()
    text_content = data.get('text', '').strip()
    transform_type = data.get('transform_type', 'humanized')
    
    if not text_content:
        return jsonify({'error': 'Text content is required'}), 400
    
    if transform_type == 'humanized':
        transformed = humanize_text(text_content)
    elif transform_type == 'ai_generated':
        transformed = formalize_text(text_content)
    else:
        transformed = text_content
    
    return jsonify({
        'transformed_text': transformed,
        'transform_type': transform_type,
    }), 200


@app.route('/api/files/<int:file_id>/download', methods=['GET'])
@jwt_required()
def download_file(file_id):
    user_id = int(get_jwt_identity())
    file = File.query.filter_by(id=file_id, user_id=user_id).first()
    if not file:
        return jsonify({'error': 'File not found'}), 404
    
    return send_from_directory(Config.UPLOAD_FOLDER, file.stored_name, 
                               as_attachment=True, download_name=file.original_name)


# ── Dashboard Endpoints ────────────────────────────────────────────────────────

@app.route('/api/dashboard/stats', methods=['GET'])
@jwt_required()
def dashboard_stats():
    user_id = int(get_jwt_identity())
    
    total = File.query.filter_by(user_id=user_id).count()
    uploaded = File.query.filter_by(user_id=user_id, status='uploaded').count()
    processing = File.query.filter_by(user_id=user_id, status='processing').count()
    completed = File.query.filter_by(user_id=user_id, status='completed').count()
    failed = File.query.filter_by(user_id=user_id, status='failed').count()
    
    recent_files = File.query.filter_by(user_id=user_id)\
        .order_by(File.created_at.desc()).limit(10).all()
    
    login_logs = LoginLog.query.filter_by(user_id=user_id)\
        .order_by(LoginLog.timestamp.desc()).limit(5).all()
    
    return jsonify({
        'stats': {
            'total_files': total,
            'uploaded': uploaded,
            'processing': processing,
            'completed': completed,
            'failed': failed,
        },
        'recent_files': [f.to_dict() for f in recent_files],
        'recent_activity': [l.to_dict() for l in login_logs],
    }), 200


# ══════════════════════════════════════════════════════════════
#  HANDWRITING STYLE ANALYSIS ENGINE
# ══════════════════════════════════════════════════════════════

# Google Fonts available in the app, categorized
_FONT_PROFILES = {
    'Caveat':              {'cursive': 0.3, 'formal': 0.1, 'messy': 0.5, 'thick': 0.4},
    'Indie Flower':        {'cursive': 0.2, 'formal': 0.0, 'messy': 0.6, 'thick': 0.3},
    'Dancing Script':      {'cursive': 0.9, 'formal': 0.5, 'messy': 0.1, 'thick': 0.4},
    'Patrick Hand':        {'cursive': 0.1, 'formal': 0.6, 'messy': 0.1, 'thick': 0.5},
    'Shadows Into Light':  {'cursive': 0.4, 'formal': 0.1, 'messy': 0.5, 'thick': 0.2},
    'Kalam':               {'cursive': 0.5, 'formal': 0.3, 'messy': 0.3, 'thick': 0.5},
    'Architects Daughter': {'cursive': 0.1, 'formal': 0.7, 'messy': 0.1, 'thick': 0.4},
    'Coming Soon':         {'cursive': 0.2, 'formal': 0.2, 'messy': 0.5, 'thick': 0.3},
    'Gochi Hand':          {'cursive': 0.3, 'formal': 0.2, 'messy': 0.4, 'thick': 0.7},
    'Handlee':             {'cursive': 0.3, 'formal': 0.5, 'messy': 0.2, 'thick': 0.4},
    'Just Another Hand':   {'cursive': 0.2, 'formal': 0.1, 'messy': 0.7, 'thick': 0.2},
    'Loved by the King':   {'cursive': 0.3, 'formal': 0.0, 'messy': 0.8, 'thick': 0.2},
    'Nothing You Could Do':{'cursive': 0.2, 'formal': 0.1, 'messy': 0.7, 'thick': 0.3},
    'Reenie Beanie':       {'cursive': 0.3, 'formal': 0.0, 'messy': 0.8, 'thick': 0.2},
    'Rock Salt':           {'cursive': 0.2, 'formal': 0.1, 'messy': 0.6, 'thick': 0.6},
    'Sacramento':          {'cursive': 0.95, 'formal': 0.7, 'messy': 0.0, 'thick': 0.3},
    'Satisfy':             {'cursive': 0.9, 'formal': 0.6, 'messy': 0.1, 'thick': 0.4},
    'Homemade Apple':      {'cursive': 0.4, 'formal': 0.2, 'messy': 0.5, 'thick': 0.6},
    'La Belle Aurore':     {'cursive': 0.7, 'formal': 0.4, 'messy': 0.3, 'thick': 0.3},
    'Cedarville Cursive':  {'cursive': 0.8, 'formal': 0.3, 'messy': 0.2, 'thick': 0.3},
}


def _analyze_handwriting_style(image_path):
    """
    Analyze a handwriting sample image and extract style characteristics.
    
    Uses PIL + numpy image processing to measure:
    - Ink color (dominant dark color)
    - Stroke weight (thickness of strokes)
    - Slant angle (lean of characters)
    - Spacing characteristics (letter/word gaps)
    - Baseline variation (how much text wobbles)
    - Connectedness (cursive vs print)
    - Character size
    
    Returns dict with font_match, ink_color, and params.
    """
    from PIL import Image, ImageOps, ImageFilter, ImageStat
    import numpy as np
    import json
    
    try:
        img = Image.open(image_path)
    except Exception as e:
        print(f"Failed to open image: {e}")
        return _default_style()
    
    # Convert to RGB if needed
    if img.mode != 'RGB':
        img = img.convert('RGB')
    
    # Resize for consistent analysis (max 2000px wide)
    w, h = img.size
    if w > 2000:
        scale = 2000 / w
        img = img.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
    
    img_array = np.array(img)
    
    # ── 1. Detect ink color ──────────────────────────────────
    gray = np.array(ImageOps.grayscale(img))
    # Find the darkest 10% of pixels (ink regions)
    threshold = np.percentile(gray, 15)
    ink_mask = gray < threshold
    
    if ink_mask.any():
        ink_pixels = img_array[ink_mask]
        avg_r = int(np.mean(ink_pixels[:, 0]))
        avg_g = int(np.mean(ink_pixels[:, 1]))
        avg_b = int(np.mean(ink_pixels[:, 2]))
        ink_color = f'#{avg_r:02x}{avg_g:02x}{avg_b:02x}'
        
        # Determine if it's blue, black, or brown ink
        if avg_b > avg_r + 30 and avg_b > avg_g + 30:
            ink_type = 'blue'
        elif avg_r > avg_b + 20 and avg_g < avg_r - 10:
            ink_type = 'brown'
        else:
            ink_type = 'black'
    else:
        ink_color = '#1a1a2e'
        ink_type = 'black'
    
    print(f"  Ink color: {ink_color} ({ink_type})")
    
    # ── 2. Binarize for structure analysis ────────────────────
    # Adaptive threshold
    bg_val = np.percentile(gray, 85)
    bin_threshold = int(bg_val * 0.65)
    binary = (gray < bin_threshold).astype(np.uint8)
    
    # ── 3. Measure stroke weight ─────────────────────────────
    # Count ink pixels in horizontal strips to estimate average stroke width
    row_sums = np.sum(binary, axis=1)
    ink_rows = row_sums[row_sums > 5]  # rows with some ink
    if len(ink_rows) > 0:
        # Average ink density per row
        avg_density = np.mean(ink_rows) / binary.shape[1]
        stroke_weight = 0.8 + avg_density * 8  # Map to 0.8 - 2.0 range
        stroke_weight = min(2.0, max(0.8, stroke_weight))
    else:
        stroke_weight = 1.0
    
    print(f"  Stroke weight: {stroke_weight:.2f}")
    
    # ── 4. Measure slant angle ───────────────────────────────
    # Analyze vertical projection at different angles
    from PIL import Image as PILImage
    gray_pil = ImageOps.grayscale(img)
    best_angle = 0
    best_score = 0
    
    for angle in range(-15, 16, 1):
        rotated = gray_pil.rotate(angle, expand=False, fillcolor=255)
        rot_arr = np.array(rotated)
        col_sums = np.sum(rot_arr < bin_threshold, axis=0)
        # Score = variance of column sums (higher = better alignment)
        score = np.var(col_sums[col_sums > 0]) if np.any(col_sums > 0) else 0
        if score > best_score:
            best_score = score
            best_angle = angle
    
    slant = -best_angle  # Positive = leaning right
    slant = max(-8, min(8, slant))  # Clamp to reasonable range
    print(f"  Slant: {slant}°")
    
    # ── 5. Find text lines and measure baseline variation ────
    row_profile = np.sum(binary, axis=1)
    # Smooth the profile
    kernel_size = max(3, binary.shape[0] // 100)
    smoothed = np.convolve(row_profile, np.ones(kernel_size)/kernel_size, mode='same')
    
    # Find line centers (peaks in the profile)
    line_threshold = np.max(smoothed) * 0.15
    in_line = smoothed > line_threshold
    line_centers = []
    start = None
    for i in range(len(in_line)):
        if in_line[i] and start is None:
            start = i
        elif not in_line[i] and start is not None:
            line_centers.append((start + i) // 2)
            start = None
    if start is not None:
        line_centers.append((start + len(in_line)) // 2)
    
    # Line height and baseline variation
    if len(line_centers) >= 2:
        line_gaps = [line_centers[i+1] - line_centers[i] for i in range(len(line_centers)-1)]
        avg_line_height = np.mean(line_gaps)
        baseline_var = np.std(line_gaps) / avg_line_height if avg_line_height > 0 else 0
        baseline_shift = min(2.0, baseline_var * 5)
    else:
        avg_line_height = binary.shape[0] * 0.1
        baseline_shift = 0.5
    
    print(f"  Lines found: {len(line_centers)}, baseline_shift: {baseline_shift:.2f}")
    
    # ── 6. Measure connectedness (cursive vs print) ──────────
    # Use connected component analysis via simple flood-fill approach
    from PIL import ImageFilter as IF
    binary_pil = Image.fromarray((1 - binary) * 255)
    # Dilate slightly to connect nearby strokes
    dilated = binary_pil.filter(IF.MinFilter(3))
    dilated_arr = np.array(dilated) < 128
    
    # Count connected components (simple row-based estimation)
    # More components per line = more printed (disconnected); fewer = more cursive
    components_per_line = []
    for center in line_centers[:10]:  # Analyze up to 10 lines
        row_start = max(0, center - int(avg_line_height * 0.4))
        row_end = min(binary.shape[0], center + int(avg_line_height * 0.4))
        line_strip = dilated_arr[row_start:row_end, :]
        col_any = np.any(line_strip, axis=0).astype(int)
        # Count transitions from 0 to 1 (start of component)
        transitions = np.sum(np.diff(col_any) == 1)
        if transitions > 0:
            components_per_line.append(transitions)
    
    if components_per_line:
        avg_components = np.mean(components_per_line)
        # More components = more printed. Typical: 5-15 for cursive, 20-40+ for print
        connectedness = max(0.0, min(0.95, 1.0 - (avg_components - 5) / 40))
    else:
        connectedness = 0.5
    
    print(f"  Connectedness: {connectedness:.2f} (avg components/line: {np.mean(components_per_line) if components_per_line else 'N/A'})")
    
    # ── 7. Measure character size ─────────────────────────────
    if len(line_centers) >= 1:
        char_height_px = avg_line_height * 0.5
        # Map pixel height to font size (rough mapping)
        font_size = max(16, min(36, int(char_height_px * 0.4)))
    else:
        font_size = 22
    
    # ── 8. Compute remaining parameters ──────────────────────
    # Rotation variation: based on baseline variation
    rotation = min(3.5, baseline_shift * 2.0 + 0.3)
    
    # Size variation: based on stroke weight consistency
    size_variation = min(0.10, baseline_shift * 0.04 + 0.01)
    
    # X-drift: based on connectedness (more connected = less drift)
    x_drift = max(0.05, (1 - connectedness) * 0.8)
    
    # Pressure: based on stroke weight variation
    pressure = min(0.20, (stroke_weight - 0.8) * 0.15 + 0.03)
    
    # Rhythm: based on baseline variation
    rhythm = min(0.30, baseline_shift * 0.15 + 0.03)
    
    # Letter gap: based on connectedness
    letter_gap = round(-connectedness * 2 + 1)
    
    # Line height for frontend
    line_height_css = max(1.5, min(3.0, 2.0 + baseline_shift * 0.3))
    
    # ── 9. Match closest Google Font ─────────────────────────
    sample_features = {
        'cursive': connectedness,
        'formal': 0.3 + (1 - baseline_shift) * 0.4,  # Less wobble = more formal
        'messy': min(1.0, baseline_shift * 0.5 + rotation * 0.1),
        'thick': (stroke_weight - 0.8) / 1.2,
    }
    
    # Also bias by ink type
    if ink_type == 'blue':
        # Blue ink → favor cursive/calligraphic fonts
        sample_features['cursive'] = min(1.0, sample_features['cursive'] + 0.2)
        sample_features['formal'] = min(1.0, sample_features['formal'] + 0.1)
    
    best_font = 'Caveat'
    best_score = float('inf')
    for font_name, profile in _FONT_PROFILES.items():
        score = sum((sample_features[k] - profile[k]) ** 2 for k in sample_features)
        if score < best_score:
            best_score = score
            best_font = font_name
    
    print(f"  Font match: {best_font} (score: {best_score:.3f})")
    
    params = {
        'rotation': round(rotation, 2),
        'baselineShift': round(baseline_shift, 2),
        'sizeVariation': round(size_variation, 3),
        'xDrift': round(x_drift, 2),
        'pressure': round(pressure, 2),
        'slant': slant,
        'connectedness': round(connectedness, 2),
        'rhythm': round(rhythm, 2),
        'strokeWeight': round(stroke_weight, 2),
        'letterGap': letter_gap,
        'fontSize': font_size,
        'lineHeight': round(line_height_css, 1),
    }
    
    print(f"  Extracted params: {params}")
    
    return {
        'font_match': best_font,
        'ink_color': ink_color,
        'params': params,
    }


def _default_style():
    """Return a sensible default style when analysis fails."""
    return {
        'font_match': 'Caveat',
        'ink_color': '#1a1a2e',
        'params': {
            'rotation': 1.5, 'baselineShift': 0.8, 'sizeVariation': 0.04,
            'xDrift': 0.3, 'pressure': 0.06, 'slant': 3,
            'connectedness': 0.5, 'rhythm': 0.10, 'strokeWeight': 1.0,
            'letterGap': 0, 'fontSize': 22, 'lineHeight': 2.0,
        },
    }


# ══════════════════════════════════════════════════════════════
#  HANDWRITING STYLE API ENDPOINTS
# ══════════════════════════════════════════════════════════════

@app.route('/api/styles/upload', methods=['POST'])
@jwt_required()
def upload_style():
    """Upload a handwriting sample image, analyze it, and save the style."""
    user_id = get_jwt_identity()
    
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    if not file or not file.filename:
        return jsonify({'error': 'No file selected'}), 400
    
    # Check file type
    ext = file.filename.rsplit('.', 1)[-1].lower() if '.' in file.filename else ''
    allowed = {'jpg', 'jpeg', 'png', 'bmp', 'tiff', 'webp', 'pdf'}
    if ext not in allowed:
        return jsonify({'error': f'File type .{ext} not supported. Use JPG, PNG, or PDF.'}), 400
    
    # Save the file
    import uuid
    stored_name = f"{uuid.uuid4().hex}.{ext}"
    styles_dir = os.path.join(Config.UPLOAD_FOLDER, 'styles')
    os.makedirs(styles_dir, exist_ok=True)
    filepath = os.path.join(styles_dir, stored_name)
    file.save(filepath)
    
    # If PDF, convert first page to image
    analysis_path = filepath
    if ext == 'pdf':
        try:
            from pdf2image import convert_from_path
            poppler_path = _find_poppler_path()
            convert_kwargs = {'dpi': 300, 'fmt': 'png', 'first_page': 1, 'last_page': 1}
            if poppler_path:
                convert_kwargs['poppler_path'] = poppler_path
            pages = convert_from_path(filepath, **convert_kwargs)
            if pages:
                img_path = filepath.rsplit('.', 1)[0] + '_page1.png'
                pages[0].save(img_path, 'PNG')
                analysis_path = img_path
        except Exception as e:
            print(f"PDF to image conversion failed: {e}")
    
    # Analyze the handwriting
    print(f"Analyzing handwriting style from: {analysis_path}")
    result = _analyze_handwriting_style(analysis_path)
    
    # Get style name
    import json
    style_name = request.form.get('name', '').strip()
    if not style_name:
        # Auto-generate name
        count = HandwritingStyle.query.filter_by(user_id=user_id).count()
        style_name = f"My Style #{count + 1}"
    
    # Save to database
    style = HandwritingStyle(
        user_id=user_id,
        name=style_name,
        stored_name=stored_name,
        font_match=result['font_match'],
        ink_color=result['ink_color'],
        params=json.dumps(result['params']),
    )
    db.session.add(style)
    db.session.commit()
    
    return jsonify({
        'message': 'Style analyzed and saved!',
        'style': style.to_dict(),
    }), 201


@app.route('/api/styles', methods=['GET'])
@jwt_required()
def list_styles():
    """List all handwriting styles for the current user."""
    user_id = get_jwt_identity()
    styles = HandwritingStyle.query.filter_by(user_id=user_id)\
        .order_by(HandwritingStyle.created_at.desc()).all()
    return jsonify({'styles': [s.to_dict() for s in styles]})


@app.route('/api/styles/<int:style_id>', methods=['GET'])
@jwt_required()
def get_style(style_id):
    """Get a single handwriting style."""
    user_id = get_jwt_identity()
    style = HandwritingStyle.query.filter_by(id=style_id, user_id=user_id).first()
    if not style:
        return jsonify({'error': 'Style not found'}), 404
    return jsonify({'style': style.to_dict()})


@app.route('/api/styles/<int:style_id>', methods=['PUT'])
@jwt_required()
def rename_style(style_id):
    """Rename a handwriting style."""
    user_id = get_jwt_identity()
    style = HandwritingStyle.query.filter_by(id=style_id, user_id=user_id).first()
    if not style:
        return jsonify({'error': 'Style not found'}), 404
    
    data = request.get_json()
    name = data.get('name', '').strip()
    if name:
        style.name = name
        db.session.commit()
    
    return jsonify({'style': style.to_dict()})


@app.route('/api/styles/<int:style_id>', methods=['DELETE'])
@jwt_required()
def delete_style(style_id):
    """Delete a handwriting style and its image."""
    user_id = get_jwt_identity()
    style = HandwritingStyle.query.filter_by(id=style_id, user_id=user_id).first()
    if not style:
        return jsonify({'error': 'Style not found'}), 404
    
    # Delete the image file
    styles_dir = os.path.join(Config.UPLOAD_FOLDER, 'styles')
    filepath = os.path.join(styles_dir, style.stored_name)
    if os.path.exists(filepath):
        os.remove(filepath)
    # Also remove converted page image if exists
    page_img = filepath.rsplit('.', 1)[0] + '_page1.png'
    if os.path.exists(page_img):
        os.remove(page_img)
    
    db.session.delete(style)
    db.session.commit()
    
    return jsonify({'message': 'Style deleted'})


@app.route('/api/styles/<int:style_id>/image', methods=['GET'])
@jwt_required()
def get_style_image(style_id):
    """Serve the handwriting sample image."""
    user_id = get_jwt_identity()
    style = HandwritingStyle.query.filter_by(id=style_id, user_id=user_id).first()
    if not style:
        return jsonify({'error': 'Style not found'}), 404
    
    styles_dir = os.path.join(Config.UPLOAD_FOLDER, 'styles')
    return send_from_directory(styles_dir, style.stored_name)


# ── Chat Endpoints ─────────────────────────────────────────────────────────────

@app.route('/api/chat/users', methods=['GET'])
@jwt_required()
def search_users():
    user_id = int(get_jwt_identity())
    query = request.args.get('q', '').strip()
    
    if query:
        users = User.query.filter(
            User.id != user_id,
            (User.name.ilike(f'%{query}%') | User.username.ilike(f'%{query}%'))
        ).limit(20).all()
    else:
        users = User.query.filter(User.id != user_id).limit(20).all()
    
    return jsonify({
        'users': [{'id': u.id, 'name': u.name, 'username': u.username, 'profile_photo': u.profile_photo} for u in users]
    }), 200


@app.route('/api/chat/rooms', methods=['GET'])
@jwt_required()
def get_chat_rooms():
    user_id = int(get_jwt_identity())
    
    memberships = ChatMember.query.filter_by(user_id=user_id).all()
    room_ids = [m.room_id for m in memberships]
    rooms = ChatRoom.query.filter(ChatRoom.id.in_(room_ids)).all() if room_ids else []
    
    rooms_data = []
    for room in rooms:
        last_msg = Message.query.filter_by(room_id=room.id)\
            .order_by(Message.created_at.desc()).first()
        
        room_dict = room.to_dict()
        room_dict['last_message'] = last_msg.to_dict() if last_msg else None
        
        # For individual chats, show the other person's name
        if room.room_type == 'individual':
            other_member = ChatMember.query.filter(
                ChatMember.room_id == room.id, ChatMember.user_id != user_id
            ).first()
            if other_member and other_member.user:
                room_dict['display_name'] = other_member.user.name
                room_dict['display_photo'] = other_member.user.profile_photo
            else:
                room_dict['display_name'] = room.name or 'Unknown'
                room_dict['display_photo'] = None
        else:
            room_dict['display_name'] = room.name
            room_dict['display_photo'] = None
        
        rooms_data.append(room_dict)
    
    return jsonify({'rooms': rooms_data}), 200


@app.route('/api/chat/rooms', methods=['POST'])
@jwt_required()
def create_chat_room():
    user_id = int(get_jwt_identity())
    data = request.get_json()
    
    room_type = data.get('type', 'individual')
    name = data.get('name', '')
    member_ids = data.get('members', [])
    
    if room_type == 'individual' and len(member_ids) == 1:
        other_user_id = member_ids[0]
        # Check if room already exists
        existing_rooms = db.session.query(ChatRoom).join(ChatMember).filter(
            ChatRoom.room_type == 'individual',
            ChatMember.user_id == user_id
        ).all()
        
        for r in existing_rooms:
            other = ChatMember.query.filter(
                ChatMember.room_id == r.id, ChatMember.user_id == other_user_id
            ).first()
            if other:
                return jsonify({'room': r.to_dict(), 'message': 'Room already exists'}), 200
    
    new_room = ChatRoom(
        name=name if room_type == 'group' else None,
        room_type=room_type,
        created_by=user_id
    )
    db.session.add(new_room)
    db.session.flush()
    
    # Add creator as member
    db.session.add(ChatMember(room_id=new_room.id, user_id=user_id))
    
    # Add other members
    for mid in member_ids:
        if mid != user_id:
            db.session.add(ChatMember(room_id=new_room.id, user_id=mid))
    
    db.session.commit()
    
    return jsonify({
        'room': new_room.to_dict(),
        'message': 'Chat room created'
    }), 201


@app.route('/api/chat/rooms/<int:room_id>/messages', methods=['GET'])
@jwt_required()
def get_messages(room_id):
    user_id = int(get_jwt_identity())
    
    # Verify membership
    member = ChatMember.query.filter_by(room_id=room_id, user_id=user_id).first()
    if not member:
        return jsonify({'error': 'Not a member of this room'}), 403
    
    page = request.args.get('page', 1, type=int)
    per_page = request.args.get('per_page', 50, type=int)
    
    messages = Message.query.filter_by(room_id=room_id)\
        .order_by(Message.created_at.desc())\
        .paginate(page=page, per_page=per_page, error_out=False)
    
    return jsonify({
        'messages': [m.to_dict() for m in reversed(messages.items)],
        'has_more': messages.has_next,
        'total': messages.total,
    }), 200


@app.route('/api/chat/rooms/<int:room_id>/share-file', methods=['POST'])
@jwt_required()
def share_file_in_chat(room_id):
    user_id = int(get_jwt_identity())
    data = request.get_json()
    file_id = data.get('file_id')
    
    member = ChatMember.query.filter_by(room_id=room_id, user_id=user_id).first()
    if not member:
        return jsonify({'error': 'Not a member of this room'}), 403
    
    file = File.query.filter_by(id=file_id, user_id=user_id).first()
    if not file:
        return jsonify({'error': 'File not found'}), 404
    
    msg = Message(
        room_id=room_id,
        sender_id=user_id,
        content=f'Shared file: {file.original_name}',
        message_type='file',
        file_id=file_id
    )
    db.session.add(msg)
    db.session.commit()
    
    socketio.emit('new_message', msg.to_dict(), room=f'room_{room_id}')
    
    return jsonify({'message': 'File shared', 'msg': msg.to_dict()}), 200


# ── Settings Endpoints ─────────────────────────────────────────────────────────

@app.route('/api/settings', methods=['GET'])
@jwt_required()
def get_settings():
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    
    return jsonify({
        'settings': {
            'theme': user.theme,
            'notifications_enabled': user.notifications_enabled,
            'email': user.email,
            'phone': user.phone,
        }
    }), 200


@app.route('/api/settings', methods=['PUT'])
@jwt_required()
def update_settings():
    user_id = int(get_jwt_identity())
    user = User.query.get(user_id)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    
    data = request.get_json()
    
    if 'theme' in data:
        user.theme = data['theme']
    if 'notifications_enabled' in data:
        user.notifications_enabled = data['notifications_enabled']
    if 'email' in data and data['email'] != user.email:
        if User.query.filter(User.email == data['email'], User.id != user_id).first():
            return jsonify({'error': 'Email already in use'}), 409
        user.email = data['email']
    if 'phone' in data:
        user.phone = data['phone']
    
    db.session.commit()
    return jsonify({'message': 'Settings updated', 'settings': {
        'theme': user.theme,
        'notifications_enabled': user.notifications_enabled,
    }}), 200


@app.route('/api/settings/delete-account', methods=['DELETE'])
@jwt_required()
def delete_account():
    user_id = int(get_jwt_identity())
    data = request.get_json()
    password = data.get('password', '')
    
    user = User.query.get(user_id)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    
    if not check_password_hash(user.password_hash, password):
        return jsonify({'error': 'Invalid password'}), 401
    
    # Delete user's files
    files = File.query.filter_by(user_id=user_id).all()
    for f in files:
        filepath = os.path.join(Config.UPLOAD_FOLDER, f.stored_name)
        if os.path.exists(filepath):
            os.remove(filepath)
    
    db.session.delete(user)
    db.session.commit()
    
    return jsonify({'message': 'Account deleted successfully'}), 200


# ── Socket.IO Events ──────────────────────────────────────────────────────────

@socketio.on('join_room')
def handle_join_room(data):
    room_id = data.get('room_id')
    user_id = data.get('user_id')
    join_room(f'room_{room_id}')
    emit('user_joined', {'user_id': user_id, 'room_id': room_id}, room=f'room_{room_id}')


@socketio.on('leave_room')
def handle_leave_room(data):
    room_id = data.get('room_id')
    user_id = data.get('user_id')
    leave_room(f'room_{room_id}')
    emit('user_left', {'user_id': user_id, 'room_id': room_id}, room=f'room_{room_id}')


@socketio.on('send_message')
def handle_send_message(data):
    room_id = data.get('room_id')
    sender_id = data.get('sender_id')
    content = data.get('content', '')
    message_type = data.get('message_type', 'text')
    file_id = data.get('file_id')
    
    with app.app_context():
        msg = Message(
            room_id=room_id,
            sender_id=sender_id,
            content=content,
            message_type=message_type,
            file_id=file_id
        )
        db.session.add(msg)
        db.session.commit()
        
        emit('new_message', msg.to_dict(), room=f'room_{room_id}')


@socketio.on('typing')
def handle_typing(data):
    room_id = data.get('room_id')
    user_name = data.get('user_name')
    emit('user_typing', {'user_name': user_name}, room=f'room_{room_id}', include_self=False)


@socketio.on('stop_typing')
def handle_stop_typing(data):
    room_id = data.get('room_id')
    emit('user_stop_typing', {}, room=f'room_{room_id}', include_self=False)


# ── Run ────────────────────────────────────────────────────────────────────────

if __name__ == '__main__':
    print("\n" + "="*60)
    print("  [*] Handwriting Transformation Backend")
    print("  [>] Running on http://127.0.0.1:5000")
    print("  [#] Upload folder:", Config.UPLOAD_FOLDER)
    print("="*60 + "\n")
    socketio.run(app, debug=True, host='0.0.0.0', port=5000, allow_unsafe_werkzeug=True)

